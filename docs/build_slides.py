# -*- coding: utf-8 -*-
"""Generate the pitch + demo slides (.pptx) for My Wallet."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Palette (amber finance brand on charcoal) ----
CHARCOAL = RGBColor(0x1E, 0x1E, 0x1E)
DARK     = RGBColor(0x26, 0x26, 0x26)
AMBER    = RGBColor(0xFF, 0xC1, 0x07)
AMBER_DK = RGBColor(0xE5, 0xA6, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
INK      = RGBColor(0x22, 0x22, 0x22)
GREY     = RGBColor(0x6B, 0x6B, 0x6B)
LGREY    = RGBColor(0xF3, 0xF3, 0xF3)
MGREY    = RGBColor(0xD9, 0xD9, 0xD9)
GREEN    = RGBColor(0x2E, 0x9E, 0x5B)
RED      = RGBColor(0xE5, 0x53, 0x3C)
CREAMTXT = RGBColor(0xBD, 0xBD, 0xBD)

HEAD_FONT = "Trebuchet MS"
BODY_FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    sp = r._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def _set_runs(p, segments, size, bold, color, font, italic, spacing):
    if isinstance(segments, str):
        segments = [(segments, color, bold)]
    for j, seg in enumerate(segments):
        txt, col, bd = seg
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bd
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = col


def text(s, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        p.space_after = Pt(ln.get("sa", 4))
        p.space_before = Pt(ln.get("sb", 0))
        if "line" in ln:
            p.line_spacing = ln["line"]
        _set_runs(
            p, ln["t"], ln.get("size", 16), ln.get("bold", False),
            ln.get("color", INK), ln.get("font", BODY_FONT),
            ln.get("italic", False), ln.get("sa", 4),
        )
    return tb


def card(s, l, t, w, h, fill, line=None, lw=1.0, radius=0.06, shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp


def rect(s, l, t, w, h, fill, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is not None:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def badge(s, l, t, d, fill, glyph, gcolor=WHITE, gsize=18, shape=MSO_SHAPE.OVAL):
    shp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = glyph
    r.font.size = Pt(gsize); r.font.bold = True
    r.font.color.rgb = gcolor; r.font.name = HEAD_FONT
    return shp


def title(s, txt, color=INK, top=0.55, l=0.85):
    # amber square motif + title (no underline)
    badge(s, l, top + 0.06, 0.26, AMBER, "", shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, l + 0.45, top, 11.5, 0.8,
         [{"t": txt, "size": 30, "bold": True, "color": color, "font": HEAD_FONT}])


def phone_placeholder(s, l, t, w, h, label, dark=False):
    fill = RGBColor(0x2C, 0x2C, 0x2C) if dark else LGREY
    bord = AMBER if dark else MGREY
    card(s, l, t, w, h, fill, line=bord, lw=1.5, radius=0.08)
    tcol = AMBER if dark else GREY
    text(s, l, t, w, h,
         [{"t": "🖼", "size": 30, "align": PP_ALIGN.CENTER, "color": tcol, "sa": 6},
          {"t": label, "size": 13, "bold": True, "align": PP_ALIGN.CENTER,
           "color": (WHITE if dark else INK), "sa": 3},
          {"t": "[ insert screenshot ]", "size": 10, "italic": True,
           "align": PP_ALIGN.CENTER, "color": (CREAMTXT if dark else GREY)}],
         anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = slide(CHARCOAL)
# left text block
text(s, 0.9, 1.9, 7.4, 1.2,
     [{"t": "My Wallet", "size": 60, "bold": True, "color": AMBER, "font": HEAD_FONT}])
text(s, 0.95, 3.15, 7.4, 0.7,
     [{"t": "Personal Expense Tracker", "size": 28, "bold": True, "color": WHITE,
       "font": HEAD_FONT}])
text(s, 0.95, 3.95, 7.4, 0.6,
     [{"t": "A cross-platform CRUD mobile app built with Flutter & Firebase",
       "size": 16, "color": CREAMTXT, "italic": True}])
# bottom meta
rect(s, 0.95, 5.35, 0.5, 0.06, AMBER)
text(s, 0.95, 5.55, 7.6, 1.2,
     [{"t": "DES3113 — Mobile App Design & Development  ·  Pitch & Demo",
       "size": 14, "color": WHITE, "sa": 4},
      {"t": "Team of 5  ·  Lecturer: Dr. Ahmad Wiraputra bin Selamat  ·  2026",
       "size": 13, "color": CREAMTXT}])
# right phone mockup
phone_placeholder(s, 9.0, 1.5, 3.3, 4.6, "Ledger / Home", dark=True)

# ============================================================
# SLIDE 2 — THE PROBLEM
# ============================================================
s = slide(WHITE)
title(s, "The Problem")
text(s, 0.9, 1.6, 6.2, 2.0,
     [{"t": "Where does my money actually go?", "size": 20, "bold": True,
       "color": INK, "font": HEAD_FONT, "sa": 10},
      {"t": "Small, frequent purchases are easy to forget and quickly add up. "
            "Most finance apps are cluttered, ad-heavy, or hide the basics "
            "behind a paywall.",
       "size": 15, "color": GREY, "line": 1.2}])
probs = [
    ("🧾", "Hard to track", "Daily cash spending slips through the cracks."),
    ("🧩", "Too complex", "Existing apps are bloated and confusing."),
    ("💾", "No backup", "Records trapped on one phone — lost if it's gone."),
    ("📉", "No insight", "Spending patterns and budgets aren't obvious."),
]
y = 1.7
for i, (ic, h, d) in enumerate(probs):
    yy = 1.65 + i * 1.18
    badge(s, 7.4, yy, 0.62, LGREY, ic, gcolor=INK, gsize=18)
    text(s, 8.2, yy - 0.02, 4.3, 1.1,
         [{"t": h, "size": 15, "bold": True, "color": INK, "sa": 2},
          {"t": d, "size": 12.5, "color": GREY}])

# ============================================================
# SLIDE 3 — THE SOLUTION
# ============================================================
s = slide(WHITE)
title(s, "The Solution")
text(s, 0.9, 1.55, 7.0, 2.4,
     [{"t": "A fast, simple, cloud-backed ledger", "size": 20, "bold": True,
       "color": INK, "font": HEAD_FONT, "sa": 10},
      {"t": "My Wallet lets you record income and expenses in seconds, browse "
            "them by day and month, see spending trends, and stay on budget — "
            "with every record stored safely in the cloud and synced in realtime.",
       "size": 15, "color": GREY, "line": 1.25, "sa": 12}])
points = [
    "Record a transaction in just a few taps",
    "Browse by day, switch between months",
    "Charts and category breakdown",
    "Monthly budget tracking",
    "Realtime cloud sync with Firebase",
]
for i, pt in enumerate(points):
    yy = 3.55 + i * 0.62
    badge(s, 0.95, yy, 0.34, AMBER, "✓", gcolor=INK, gsize=13)
    text(s, 1.45, yy + 0.0, 6.4, 0.5,
         [{"t": pt, "size": 14.5, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
phone_placeholder(s, 9.0, 1.55, 3.3, 4.7, "My Wallet — Ledger")

# ============================================================
# SLIDE 4 — KEY FEATURES (grid)
# ============================================================
s = slide(WHITE)
title(s, "Key Features")
feats = [
    ("🧾", "Smart Ledger", "Transactions grouped by day with monthly totals."),
    ("➕", "Quick Record", "Custom keypad, category grid, date & note."),
    ("📊", "Charts", "Week / Month / Year trend + category breakdown."),
    ("🎯", "Budget", "Set a monthly limit; track remaining balance."),
    ("👤", "Statistics", "Lifetime income, expense, balance & activity."),
    ("☁️", "Cloud Database", "Stored in Cloud Firestore, synced in realtime."),
]
cw, ch, gx, gy = 3.74, 1.62, 0.3, 0.3
x0, y0 = 0.9, 1.7
for i, (ic, h, d) in enumerate(feats):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    card(s, x, y, cw, ch, LGREY, radius=0.07)
    badge(s, x + 0.25, y + 0.28, 0.7, AMBER, ic, gcolor=INK, gsize=20)
    text(s, x + 1.1, y + 0.22, cw - 1.3, ch - 0.3,
         [{"t": h, "size": 15, "bold": True, "color": INK, "sa": 3},
          {"t": d, "size": 11.5, "color": GREY, "line": 1.1}])

# ============================================================
# SLIDE 5 — CRUD AT THE CORE
# ============================================================
s = slide(WHITE)
title(s, "CRUD at the Core")
text(s, 0.9, 1.45, 11.5, 0.5,
     [{"t": "Every requirement maps to a real feature, all driven by one Provider.",
       "size": 14, "color": GREY, "italic": True}])
crud = [
    ("C", "Create", AMBER, "Add a record via the keypad, category grid, date picker and note."),
    ("R", "Read", GREEN, "Browse the Ledger grouped by day; open any record in Detail."),
    ("U", "Update", RGBColor(0x3B,0x82,0xC4), "Edit a record — the Add screen reopens pre-filled."),
    ("D", "Delete", RED, "Remove a record with a confirm dialog, or clear all."),
]
cw, gx = 2.86, 0.24
x0, y = 0.9, 2.25
for i, (letter, name, col, d) in enumerate(crud):
    x = x0 + i * (cw + gx)
    card(s, x, y, cw, 3.4, WHITE, line=MGREY, lw=1.0, radius=0.05)
    rect(s, x, y, cw, 0.12, col)
    badge(s, x + cw/2 - 0.55, y + 0.45, 1.1, col, letter, gcolor=WHITE, gsize=40)
    text(s, x + 0.2, y + 1.75, cw - 0.4, 0.5,
         [{"t": name, "size": 19, "bold": True, "color": INK, "font": HEAD_FONT,
           "align": PP_ALIGN.CENTER}])
    text(s, x + 0.28, y + 2.35, cw - 0.56, 1.0,
         [{"t": d, "size": 12, "color": GREY, "align": PP_ALIGN.CENTER, "line": 1.15}])

# ============================================================
# SLIDE 6 — ARCHITECTURE
# ============================================================
s = slide(WHITE)
title(s, "Architecture")
text(s, 0.9, 1.5, 11.5, 0.5,
     [{"t": "A clean, layered design — the UI never touches storage directly.",
       "size": 14, "color": GREY, "italic": True}])
layers = [
    ("UI Layer", "Ledger · Charts · Discover · Profile · Add / Edit", AMBER, INK),
    ("State Layer", "TransactionProvider  (ChangeNotifier) — CRUD, totals, budget", DARK, WHITE),
    ("Model Layer", "Transaction · Category · TransactionType", RGBColor(0x55,0x55,0x55), WHITE),
    ("Storage Layer", "Cloud Firestore  (realtime NoSQL database)", RGBColor(0x80,0x80,0x80), WHITE),
]
y = 2.25
lw_ = 8.6
x = 0.9
arrows = ["context.watch / read", "toMap / fromMap", "snapshots / set / delete"]
for i, (name, desc, col, tcol) in enumerate(layers):
    card(s, x, y, lw_, 0.92, col, radius=0.05)
    text(s, x + 0.35, y, lw_ - 0.6, 0.92,
         [{"t": [(name + "   ", tcol, True), (desc, tcol, False)], "size": 14.5,
           "color": tcol}], anchor=MSO_ANCHOR.MIDDLE)
    if i < len(arrows):
        text(s, x, y + 0.92, lw_, 0.42,
             [{"t": "↓  " + arrows[i], "size": 11, "italic": True, "color": GREY,
               "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.34
# side note
card(s, 9.85, 2.25, 2.6, 4.42, LGREY, radius=0.06)
text(s, 10.1, 2.55, 2.1, 4.0,
     [{"t": "Why it works", "size": 14, "bold": True, "color": INK,
       "font": HEAD_FONT, "sa": 8},
      {"t": "• Single source of truth", "size": 12, "color": GREY, "sa": 6},
      {"t": "• UI rebuilds on change", "size": 12, "color": GREY, "sa": 6},
      {"t": "• Logic easy to test", "size": 12, "color": GREY, "sa": 6},
      {"t": "• Realtime cloud sync", "size": 12, "color": GREY, "sa": 6}])

# ============================================================
# SLIDE 7 — TECH STACK
# ============================================================
s = slide(WHITE)
title(s, "Technology Stack")
stack = [
    ("Flutter", "Cross-platform UI toolkit (Material 3)"),
    ("Dart", "Application language (SDK ^3.11)"),
    ("firebase_core", "Firebase initialisation"),
    ("cloud_firestore", "Realtime cloud database (CRUD)"),
    ("provider", "ChangeNotifier state management"),
    ("intl · uuid", "Formatting & unique IDs"),
    ("flutter_test", "Unit & widget testing"),
    ("Git + GitHub", "Version control & collaboration"),
]
cw, ch, gx, gy = 2.86, 1.32, 0.24, 0.28
x0, y0 = 0.9, 1.75
for i, (name, d) in enumerate(stack):
    r, c = divmod(i, 4)
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    card(s, x, y, cw, ch, WHITE, line=MGREY, lw=1.0, radius=0.07)
    rect(s, x, y, 0.1, ch, AMBER)
    text(s, x + 0.3, y + 0.2, cw - 0.45, ch - 0.3,
         [{"t": name, "size": 15, "bold": True, "color": INK, "sa": 4,
           "font": HEAD_FONT},
          {"t": d, "size": 11, "color": GREY, "line": 1.1}])

# ============================================================
# SLIDE 8 — STATE & PERSISTENCE (code-ish)
# ============================================================
s = slide(WHITE)
title(s, "State & Persistence")
text(s, 0.9, 1.55, 5.6, 4.6,
     [{"t": "Provider as the single source of truth", "size": 17, "bold": True,
       "color": INK, "font": HEAD_FONT, "sa": 8},
      {"t": "TransactionProvider extends ChangeNotifier. Each CRUD method mutates "
            "the list, saves to storage, then calls notifyListeners() so every "
            "watching widget rebuilds automatically.", "size": 13.5, "color": GREY,
       "line": 1.25, "sa": 12},
      {"t": "Cloud persistence", "size": 17, "bold": True, "color": INK,
       "font": HEAD_FONT, "sa": 8},
      {"t": "Each transaction is a document in Cloud Firestore. A realtime snapshot "
            "listener streams changes back, so the UI updates instantly across devices.",
       "size": 13.5, "color": GREY, "line": 1.25}])
# code card
card(s, 6.9, 1.6, 5.5, 4.5, CHARCOAL, radius=0.04)
code = [
    ("// Realtime read", CREAMTXT),
    ("_txCol.snapshots().listen((snap) {", AMBER),
    ("  _transactions = snap.docs", WHITE),
    ("      .map((d) =>", WHITE),
    ("          Transaction.fromMap(d.data()))", WHITE),
    ("      .toList();", WHITE),
    ("  notifyListeners();", WHITE),
    ("});", AMBER),
    ("", WHITE),
    ("// Create / Update", CREAMTXT),
    ("await _txCol.doc(t.id).set(t.toMap());", WHITE),
    ("", WHITE),
    ("// Delete", CREAMTXT),
    ("await _txCol.doc(id).delete();", WHITE),
]
tb = s.shapes.add_textbox(Inches(7.15), Inches(1.85), Inches(5.1), Inches(4.0))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
for i, (ln, col) in enumerate(code):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(1)
    r = p.add_run(); r.text = ln if ln else " "
    r.font.size = Pt(11.5); r.font.name = "Consolas"; r.font.color.rgb = col

# ============================================================
# SLIDE 9 — UI / UX DESIGN
# ============================================================
s = slide(WHITE)
title(s, "UI / UX Design")
text(s, 0.9, 1.5, 11.5, 0.5,
     [{"t": "Material 3 · amber brand colour · familiar daily-ledger layout.",
       "size": 14, "color": GREY, "italic": True}])
shots = ["Ledger", "Add Record", "Charts", "Profile"]
sw_, sh_, gx = 2.78, 3.9, 0.24
x0, y = 0.9, 2.15
for i, name in enumerate(shots):
    x = x0 + i * (sw_ + gx)
    phone_placeholder(s, x, y, sw_, sh_, name)

# ============================================================
# SLIDE 10 — LIVE DEMO
# ============================================================
s = slide(CHARCOAL)
text(s, 0.9, 2.35, 8.0, 1.0,
     [{"t": "Live Demo", "size": 48, "bold": True, "color": AMBER,
       "font": HEAD_FONT}])
rect(s, 0.95, 3.5, 0.6, 0.07, WHITE)
steps = [
    "Add an expense with the keypad",
    "Watch it appear in today's ledger",
    "Open it, edit the amount, delete it",
    "Set a budget · view the charts",
]
for i, st in enumerate(steps):
    yy = 3.85 + i * 0.66
    badge(s, 0.95, yy, 0.42, AMBER, str(i + 1), gcolor=INK, gsize=16)
    text(s, 1.55, yy, 7.0, 0.5,
         [{"t": st, "size": 16, "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
phone_placeholder(s, 9.4, 1.7, 3.0, 4.1, "App Demo", dark=True)

# ============================================================
# SLIDE 11 — TESTING
# ============================================================
s = slide(WHITE)
title(s, "Testing")
# big stat
card(s, 0.9, 1.75, 4.0, 2.3, CHARCOAL, radius=0.06)
text(s, 0.9, 2.0, 4.0, 1.4,
     [{"t": "12 / 12", "size": 52, "bold": True, "color": AMBER,
       "align": PP_ALIGN.CENTER, "font": HEAD_FONT}],
     anchor=MSO_ANCHOR.TOP)
text(s, 0.9, 3.25, 4.0, 0.7,
     [{"t": "automated tests passing", "size": 14, "color": WHITE,
       "align": PP_ALIGN.CENTER}])
# breakdown
text(s, 5.3, 1.75, 7.0, 0.5,
     [{"t": "What we tested", "size": 17, "bold": True, "color": INK,
       "font": HEAD_FONT}])
tests = [
    ("Unit — Model", "Serialisation round-trips, copyWith, category labels."),
    ("Unit — Provider", "CRUD + budget on an in-memory test database."),
    ("Widget", "App builds and the bottom navigation renders."),
    ("Manual", "9 flows incl. data appearing live in the Firebase console."),
]
for i, (h, d) in enumerate(tests):
    yy = 2.35 + i * 0.92
    badge(s, 5.3, yy, 0.4, AMBER, "✓", gcolor=INK, gsize=14)
    text(s, 5.85, yy - 0.04, 6.6, 0.9,
         [{"t": h, "size": 14.5, "bold": True, "color": INK, "sa": 2},
          {"t": d, "size": 12, "color": GREY}])

# ============================================================
# SLIDE 12 — TEAM & CONTRIBUTIONS
# ============================================================
s = slide(WHITE)
title(s, "Team & Contributions")
team = [
    ("Sazit Ul Islam", "Project lead · State management & data model"),
    ("Nur Aisyatul Najwa binti Mohamad Nasir", "UI/UX · Ledger & Detail screens"),
    ("Nurfatin Aqilah binti Zolkifli", "Add/Edit & Charts screens"),
    ("Shi Kaiyan", "Discover, Profile · Firebase & testing"),
    ("Mst Samiya Haque Kotha", "Documentation, report & presentation"),
]
y0 = 1.8
for i, (name, role) in enumerate(team):
    yy = y0 + i * 0.92
    card(s, 0.9, yy, 11.5, 0.78, LGREY if i % 2 == 0 else WHITE,
         line=(None if i % 2 == 0 else MGREY), lw=1.0, radius=0.08)
    badge(s, 1.15, yy + 0.13, 0.52, AMBER, str(i + 1), gcolor=INK, gsize=18)
    text(s, 1.95, yy, 5.1, 0.78,
         [{"t": name, "size": 12.5, "bold": True, "color": INK}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 7.2, yy, 5.0, 0.78,
         [{"t": role, "size": 12, "color": GREY}], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.9, 6.55, 11.5, 0.5,
     [{"t": "Every member contributed across design, development, testing & docs. "
            "Source managed with Git + GitHub.", "size": 11.5, "italic": True,
       "color": GREY}])

# ============================================================
# SLIDE 13 — CHALLENGES & FUTURE
# ============================================================
s = slide(WHITE)
title(s, "Challenges & What's Next")
text(s, 0.9, 1.7, 5.6, 0.5,
     [{"t": "Challenges we solved", "size": 17, "bold": True, "color": INK,
       "font": HEAD_FONT}])
ch = [
    "Keeping UI in sync → single Provider",
    "Cloud data → Firestore + realtime listener",
    "Trend chart → custom CustomPainter",
    "Valid input → keypad guard rules",
]
for i, c in enumerate(ch):
    yy = 2.4 + i * 0.74
    badge(s, 0.95, yy, 0.36, AMBER, "✓", gcolor=INK, gsize=13)
    text(s, 1.45, yy, 5.0, 0.6,
         [{"t": c, "size": 13.5, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
text(s, 6.9, 1.7, 5.5, 0.5,
     [{"t": "Future enhancements", "size": 17, "bold": True, "color": INK,
       "font": HEAD_FONT}])
fut = [
    "User accounts (Firebase Auth)",
    "Firestore security rules",
    "Export to CSV / PDF",
    "Recurring bills & reminders",
    "Dark mode & multi-currency",
]
for i, f in enumerate(fut):
    yy = 2.4 + i * 0.64
    badge(s, 6.95, yy, 0.36, DARK, "→", gcolor=AMBER, gsize=13)
    text(s, 7.45, yy, 5.0, 0.6,
         [{"t": f, "size": 13.5, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# SLIDE 14 — THANK YOU
# ============================================================
s = slide(CHARCOAL)
text(s, 0.9, 2.5, 11.5, 1.2,
     [{"t": "Thank You", "size": 56, "bold": True, "color": AMBER,
       "font": HEAD_FONT}])
text(s, 0.95, 3.85, 11.5, 0.7,
     [{"t": "My Wallet — Personal Expense Tracker", "size": 22, "bold": True,
       "color": WHITE, "font": HEAD_FONT}])
text(s, 0.95, 4.6, 11.5, 0.6,
     [{"t": "Built with Flutter  ·  Firebase  ·  Cloud Firestore",
       "size": 15, "color": CREAMTXT, "italic": True}])
text(s, 0.95, 5.5, 11.5, 0.6,
     [{"t": "Questions & demo  ·  github.com/LBBT-God/expense-tracker", "size": 14,
       "color": AMBER}])

out = r"D:\des3113\docs\My_Wallet_Presentation.pptx"
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
